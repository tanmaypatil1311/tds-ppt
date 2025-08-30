from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
import tempfile
import logging
import base64
import io
from typing import Dict, List, Any, Optional
import uuid

logger = logging.getLogger(__name__)


class PPTXGenerator:
    def __init__(self):
        self.temp_dir = tempfile.mkdtemp()

   #  def generate_presentation(self,
   #                            slides: List[Dict],
   #                            template_data: Dict,
   #                            template_path: str,
   #                            options: Dict = None) -> str:
   #      """Generate a PowerPoint presentation with template styling"""

   #      if options is None:
   #          options = {}

   #      try:
   #          # Create new presentation from template
   #          if os.path.exists(template_path):
   #              prs = Presentation(template_path)
   #              # Clear existing slides if any
   #              while len(prs.slides) > 0:
   #                  r_id = prs.slides._sldIdLst[0].rId
   #                  prs.part.drop_rel(r_id)
   #                  del prs.slides._sldIdLst[0]
   #          else:
   #              # Create blank presentation
   #              prs = Presentation()

   #          # Generate slides
   #          for slide_data in slides:
   #              self._create_slide(prs, slide_data, template_data, options)

   #          # Save presentation
   #          output_path = os.path.join(
   #              self.temp_dir, f"generated_{uuid.uuid4().hex}.pptx")
   #          prs.save(output_path)

   #          logger.info(
   #              f"Presentation generated successfully: {len(slides)} slides")
   #          return output_path

   #      except Exception as e:
   #          logger.error(f"Error generating presentation: {e}")
   #          raise

    def generate_presentation(self, slides, template_data, template_path, options={}):
        # Load template if provided, otherwise start fresh
        prs = Presentation(template_path) if template_path else Presentation()

        # Remove existing empty slides (optional)
        while prs.slides:
            r_id = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(r_id)
            del prs.slides._sldIdLst[0]

        # Loop through AI-generated slide data
        for slide in slides:
            # Pick first available layout from template
            layout = prs.slide_layouts[0] if prs.slide_layouts else prs.slide_layouts[0]
            slide_obj = prs.slides.add_slide(layout)

            # Title
            if "title" in slide and slide["title"]:
                if slide_obj.shapes.title:
                    slide_obj.shapes.title.text = slide["title"]

            # Content (bullet points or text)
            if "content" in slide and slide["content"]:
                # Use first placeholder for body text
                body_shapes = [
                    s for s in slide_obj.placeholders if s.is_placeholder and s.placeholder_format.type == 1]
                if body_shapes:
                    tf = body_shapes[0].text_frame
                    tf.clear()
                    for line in slide["content"].split("\n"):
                        p = tf.add_paragraph()
                        p.text = line.strip()
                        p.font.size = Pt(18)

            # Speaker Notes
            if "notes" in slide and slide["notes"]:
                notes_slide = slide_obj.notes_slide
                text_frame = notes_slide.notes_text_frame
                text_frame.text = slide["notes"]

        # Save final presentation
        output_path = template_path.replace(
            ".pptx", "_generated.pptx") if template_path else "generated_presentation.pptx"
        prs.save(output_path)
        return output_path

    def _create_slide(self, prs: Presentation, slide_data: Dict, template_data: Dict, options: Dict):
        """Create a single slide"""

        slide_type = slide_data.get('slide_type', 'content')

        # Get appropriate layout
        layout = self._get_layout_for_slide_type(
            prs, template_data, slide_type)

        # Add slide
        slide = prs.slides.add_slide(layout)

        # Add content based on slide type
        if slide_type == 'title':
            self._populate_title_slide(slide, slide_data, template_data)
        elif slide_type == 'section':
            self._populate_section_slide(slide, slide_data, template_data)
        else:  # content, conclusion, or other
            self._populate_content_slide(
                slide, slide_data, template_data, options)

        # Add speaker notes if available
        if slide_data.get('notes'):
            self._add_speaker_notes(slide, slide_data['notes'])

    def _get_layout_for_slide_type(self, prs: Presentation, template_data: Dict, slide_type: str):
        """Get the most appropriate layout for the slide type"""

        layouts = template_data.get('layouts', [])

        # Mapping preferences
        type_preferences = {
            'title': ['title', 'basic'],
            'content': ['content', 'basic'],
            'section': ['title', 'content'],
            'conclusion': ['content', 'basic']
        }

        preferred_types = type_preferences.get(
            slide_type, ['content', 'basic'])

        # Find matching layout
        for pref_type in preferred_types:
            for layout_info in layouts:
                if layout_info.get('used_for') == pref_type:
                    layout_index = layout_info.get('index', 0)
                    if layout_index < len(prs.slide_layouts):
                        return prs.slide_layouts[layout_index]

        # Fallback to first available layout
        return prs.slide_layouts[0] if prs.slide_layouts else None

    def _populate_title_slide(self, slide, slide_data: Dict, template_data: Dict):
        """Populate a title slide"""

        title = slide_data.get('title', 'Presentation Title')
        subtitle_content = slide_data.get('content', [])
        subtitle = subtitle_content[0] if subtitle_content else ''

        # Find title and subtitle placeholders
        title_placeholder = None
        subtitle_placeholder = None

        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == 0:  # Title
                title_placeholder = placeholder
            elif placeholder.placeholder_format.type == 1:  # Subtitle
                subtitle_placeholder = placeholder

        # Set title
        if title_placeholder:
            title_placeholder.text = title
            self._apply_text_formatting(
                title_placeholder, template_data, 'title')

        # Set subtitle
        if subtitle_placeholder and subtitle:
            subtitle_placeholder.text = subtitle
            self._apply_text_formatting(
                subtitle_placeholder, template_data, 'subtitle')

    def _populate_section_slide(self, slide, slide_data: Dict, template_data: Dict):
        """Populate a section divider slide"""
        # Similar to title slide but with section-specific styling
        self._populate_title_slide(slide, slide_data, template_data)

    def _populate_content_slide(self, slide, slide_data: Dict, template_data: Dict, options: Dict):
        """Populate a content slide with title and bullet points"""

        title = slide_data.get('title', '')
        content_items = slide_data.get('content', [])

        # Find title placeholder
        title_placeholder = None
        content_placeholder = None

        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == 0:  # Title
                title_placeholder = placeholder
            # Content placeholders
            elif placeholder.placeholder_format.type in [1, 2, 7]:
                content_placeholder = placeholder

        # Set title
        if title_placeholder:
            title_placeholder.text = title
            self._apply_text_formatting(
                title_placeholder, template_data, 'title')

        # Set content
        if content_placeholder and content_items:
            self._populate_text_placeholder(
                content_placeholder, content_items, template_data)

        # Add images if available and requested
        if options.get('include_images', True):
            self._try_add_image_to_slide(slide, slide_data, template_data)

    def _populate_text_placeholder(self, placeholder, content_items: List[str], template_data: Dict):
        """Populate a text placeholder with bullet points"""

        text_frame = placeholder.text_frame
        text_frame.clear()

        for i, item in enumerate(content_items[:8]):  # Limit to 8 items
            if i == 0:
                # Use existing paragraph
                p = text_frame.paragraphs[0]
            else:
                # Add new paragraph
                p = text_frame.add_paragraph()

            p.text = str(item)
            p.level = 0

            # Apply formatting
            self._apply_paragraph_formatting(p, template_data)

    def _apply_text_formatting(self, placeholder, template_data: Dict, text_type: str = 'body'):
        """Apply formatting to text placeholder"""

        try:
            fonts = template_data.get('fonts', {})
            colors = template_data.get('colors', ['#000000'])

            text_frame = placeholder.text_frame

            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    # Set font
                    if text_type == 'title':
                        run.font.name = fonts.get('title_font', 'Calibri')
                        run.font.size = Pt(32)
                        run.font.bold = True
                    else:
                        run.font.name = fonts.get('body_font', 'Calibri')
                        run.font.size = Pt(18)

                    # Set color
                    # Avoid white text on white background
                    if colors and colors[0] != '#FFFFFF':
                        try:
                            color_hex = colors[0].replace('#', '')
                            r = int(color_hex[0:2], 16)
                            g = int(color_hex[2:4], 16)
                            b = int(color_hex[4:6], 16)
                            run.font.color.rgb = RGBColor(r, g, b)
                        except:
                            pass  # Use default color if parsing fails

        except Exception as e:
            logger.warning(f"Error applying text formatting: {e}")

    def _apply_paragraph_formatting(self, paragraph, template_data: Dict):
        """Apply formatting to a paragraph"""

        try:
            fonts = template_data.get('fonts', {})

            for run in paragraph.runs:
                run.font.name = fonts.get('body_font', 'Calibri')
                run.font.size = Pt(18)

                # Apply color if available
                colors = template_data.get('colors', [])
                if colors and len(colors) > 1:  # Use second color for body text
                    try:
                        color_hex = colors[1].replace('#', '')
                        r = int(color_hex[0:2], 16)
                        g = int(color_hex[2:4], 16)
                        b = int(color_hex[4:6], 16)
                        run.font.color.rgb = RGBColor(r, g, b)
                    except:
                        pass

        except Exception as e:
            logger.warning(f"Error applying paragraph formatting: {e}")

    def _try_add_image_to_slide(self, slide, slide_data: Dict, template_data: Dict):
        """Try to add an appropriate image to the slide"""

        try:
            images = template_data.get('images', [])
            if not images:
                return

            # Simple strategy: use first available image
            image_data = images[0]

            # Check if slide has image placeholder
            image_placeholder = None
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.type == 8:  # Picture placeholder
                    image_placeholder = placeholder
                    break

            if image_placeholder:
                # Use placeholder
                self._insert_image_in_placeholder(
                    image_placeholder, image_data)
            else:
                # Add as free-floating image in available space
                self._add_floating_image(slide, image_data)

        except Exception as e:
            logger.warning(f"Error adding image to slide: {e}")

    def _insert_image_in_placeholder(self, placeholder, image_data: Dict):
        """Insert image into a placeholder"""

        try:
            # Decode base64 image
            image_bytes = base64.b64decode(image_data['data'])

            # Create temporary file
            temp_image_path = os.path.join(
                self.temp_dir, f"temp_image_{uuid.uuid4().hex}.png")
            with open(temp_image_path, 'wb') as f:
                f.write(image_bytes)

            # Insert image
            placeholder.insert_picture(temp_image_path)

            # Clean up
            os.remove(temp_image_path)

        except Exception as e:
            logger.warning(f"Error inserting image in placeholder: {e}")

    def _add_floating_image(self, slide, image_data: Dict):
        """Add image as floating element"""

        try:
            # Decode base64 image
            image_bytes = base64.b64decode(image_data['data'])

            # Create temporary file
            temp_image_path = os.path.join(
                self.temp_dir, f"temp_image_{uuid.uuid4().hex}.png")
            with open(temp_image_path, 'wb') as f:
                f.write(image_bytes)

            # Calculate position (right side of slide)
            slide_width = slide.slide_layout.slide_master.slide_width
            slide_height = slide.slide_layout.slide_master.slide_height

            img_width = min(Inches(3), slide_width * 0.3)
            img_height = min(Inches(2), slide_height * 0.3)

            left = slide_width - img_width - Inches(0.5)
            top = Inches(1)

            # Add image
            slide.shapes.add_picture(
                temp_image_path, left, top, img_width, img_height)

            # Clean up
            os.remove(temp_image_path)

        except Exception as e:
            logger.warning(f"Error adding floating image: {e}")

    def _add_speaker_notes(self, slide, notes_text: str):
        """Add speaker notes to slide"""

        try:
            notes_slide = slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame
            notes_text_frame.text = notes_text

        except Exception as e:
            logger.warning(f"Error adding speaker notes: {e}")

    def enhance_presentation(self,
                             presentation_path: str,
                             template_data: Dict,
                             enhancements: Dict) -> str:
        """Apply enhancements to existing presentation"""

        try:
            prs = Presentation(presentation_path)

            # Apply enhancements
            if enhancements.get('improve_typography', False):
                self._improve_typography(prs, template_data)

            if enhancements.get('add_transitions', False):
                self._add_slide_transitions(prs)

            if enhancements.get('optimize_layouts', False):
                self._optimize_slide_layouts(prs, template_data)

            # Save enhanced version
            enhanced_path = presentation_path.replace(
                '.pptx', '_enhanced.pptx')
            prs.save(enhanced_path)

            return enhanced_path

        except Exception as e:
            logger.error(f"Error enhancing presentation: {e}")
            return presentation_path  # Return original on failure

    def _improve_typography(self, prs: Presentation, template_data: Dict):
        """Improve typography throughout presentation"""

        fonts = template_data.get('fonts', {})

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            # Ensure consistent font sizing
                            if run.font.size and run.font.size < Pt(14):
                                run.font.size = Pt(16)

                            # Apply template fonts
                            if not run.font.name or run.font.name == 'Calibri':
                                run.font.name = fonts.get(
                                    'body_font', 'Calibri')

    def _add_slide_transitions(self, prs: Presentation):
        """Add subtle slide transitions"""

        # This is a placeholder - python-pptx has limited transition support
        # In a full implementation, you'd need to work with the underlying XML
        logger.info("Slide transitions would be added here")

    def _optimize_slide_layouts(self, prs: Presentation, template_data: Dict):
        """Optimize slide layouts for better visual balance"""

        for slide in prs.slides:
            # Adjust text box sizes and positions for better readability
            text_shapes = [s for s in slide.shapes if hasattr(s, 'text_frame')]

            for shape in text_shapes:
                # Ensure minimum margins
                if shape.left < Inches(0.5):
                    shape.left = Inches(0.5)

                if shape.top < Inches(0.5):
                    shape.top = Inches(0.5)

                # Adjust text frame properties
                text_frame = shape.text_frame
                text_frame.margin_left = Inches(0.1)
                text_frame.margin_right = Inches(0.1)
                text_frame.margin_top = Inches(0.1)
                text_frame.margin_bottom = Inches(0.1)

    def create_presentation_preview(self, slides: List[Dict], template_data: Dict) -> List[Dict]:
        """Create preview data for slides without generating full presentation"""

        previews = []

        for slide_data in slides:
            preview = {
                'slide_number': slide_data.get('slide_number', 1),
                'title': slide_data.get('title', ''),
                # First 3 items
                'content_preview': slide_data.get('content', [])[:3],
                'slide_type': slide_data.get('slide_type', 'content'),
                'has_notes': bool(slide_data.get('notes', '')),
                'estimated_elements': self._count_slide_elements(slide_data),
                'layout_suggestion': self._suggest_layout(slide_data, template_data)
            }

            previews.append(preview)

        return previews

    def _count_slide_elements(self, slide_data: Dict) -> Dict:
        """Count elements that will be on the slide"""

        return {
            'text_blocks': len(slide_data.get('content', [])),
            'has_title': bool(slide_data.get('title', '')),
            'has_notes': bool(slide_data.get('notes', '')),
            'estimated_images': 1 if slide_data.get('slide_type') != 'title' else 0
        }

    def _suggest_layout(self, slide_data: Dict, template_data: Dict) -> str:
        """Suggest best layout for slide"""

        slide_type = slide_data.get('slide_type', 'content')
        content_count = len(slide_data.get('content', []))

        if slide_type == 'title':
            return 'Title Slide'
        elif content_count > 5:
            return 'Content Heavy'
        elif content_count <= 2:
            return 'Minimal Content'
        else:
            return 'Standard Content'
