from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
import tempfile
import logging
from typing import Dict, List, Any, Optional
import base64
from PIL import Image
import io

logger = logging.getLogger(__name__)


class PPTXAnalyzer:
    def __init__(self):
        self.supported_formats = ['.pptx', '.potx']

    def analyze_template(self, template_path: str) -> Dict[str, Any]:
        """Analyze a PowerPoint template and extract styling information"""

        if not os.path.exists(template_path):
            raise FileNotFoundError(
                f"Template file not found: {template_path}")

        try:
            prs = Presentation(template_path)

            analysis_data = {
                'layouts': self._extract_layouts(prs),
                'theme': self._extract_theme(prs),
                'colors': self._extract_colors(prs),
                'fonts': self._extract_fonts(prs),
                'images': self._extract_images(prs),
                'slide_size': self._get_slide_size(prs),
                'master_slides': self._analyze_master_slides(prs)
            }

            logger.info(
                f"Template analysis completed: {len(analysis_data['layouts'])} layouts found")
            return analysis_data

        except Exception as e:
            logger.error(f"Error analyzing template: {e}")
            raise

    def _extract_layouts(self, prs: Presentation) -> List[Dict]:
        """Extract layout information from slide master"""
        layouts = []

        try:
            for i, slide_layout in enumerate(prs.slide_layouts):
                layout_info = {
                    'index': i,
                    'name': slide_layout.name,
                    'placeholders': [],
                    'used_for': self._determine_layout_usage(slide_layout)
                }

                # Extract placeholder information
                for placeholder in slide_layout.placeholders:
                    placeholder_info = {
                        'type': placeholder.placeholder_format.type,
                        'idx': placeholder.placeholder_format.idx,
                        'left': placeholder.left,
                        'top': placeholder.top,
                        'width': placeholder.width,
                        'height': placeholder.height
                    }
                    layout_info['placeholders'].append(placeholder_info)

                layouts.append(layout_info)

        except Exception as e:
            logger.warning(f"Error extracting layouts: {e}")

        return layouts

    def _determine_layout_usage(self, slide_layout) -> str:
        """Determine the best usage for a layout based on placeholders"""
        placeholder_types = [
            p.placeholder_format.type for p in slide_layout.placeholders]

        # Common layout patterns
        if 1 in placeholder_types and 0 in placeholder_types:  # Title and content
            return 'content'
        # Title only
        elif 0 in placeholder_types and len(placeholder_types) == 1:
            return 'title'
        elif 13 in placeholder_types:  # Picture placeholder
            return 'image'
        elif len(placeholder_types) > 3:
            return 'complex'
        else:
            return 'basic'

    def _extract_theme(self, prs: Presentation) -> Dict:
        """Extract theme information"""
        theme_data = {
            'name': 'Default',
            'color_scheme': [],
            'font_scheme': {}
        }

        try:
            # Get theme from slide master
            if prs.slide_master and hasattr(prs.slide_master, 'theme'):
                theme = prs.slide_master.theme
                theme_data['name'] = getattr(theme, 'name', 'Default')
        except Exception as e:
            logger.warning(f"Error extracting theme: {e}")

        return theme_data

    def _extract_colors(self, prs: Presentation) -> List[str]:
        """Extract color palette from the presentation"""
        colors = set()

        try:
            # Extract colors from slide master
            for slide in prs.slides:
                for shape in slide.shapes:
                    # Get fill colors
                    if hasattr(shape, 'fill') and shape.fill.type is not None:
                        try:
                            if hasattr(shape.fill, 'fore_color') and hasattr(shape.fill.fore_color, 'rgb'):
                                rgb = shape.fill.fore_color.rgb
                                if rgb is not None:
                                    colors.add(f"#{rgb:06X}")
                        except:
                            pass

                    # Get text colors
                    if hasattr(shape, 'text_frame'):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                try:
                                    if hasattr(run.font, 'color') and hasattr(run.font.color, 'rgb'):
                                        rgb = run.font.color.rgb
                                        if rgb is not None:
                                            colors.add(f"#{rgb:06X}")
                                except:
                                    pass

            # Add default colors if none found
            if not colors:
                colors = {'#000000', '#FFFFFF',
                          '#1F497D', '#4F81BD', '#9CBB58'}

        except Exception as e:
            logger.warning(f"Error extracting colors: {e}")
            colors = {'#000000', '#FFFFFF', '#1F497D', '#4F81BD', '#9CBB58'}

        return list(colors)[:10]  # Limit to 10 colors

    def _extract_fonts(self, prs: Presentation) -> Dict:
        """Extract font information"""
        fonts = {
            'title_font': 'Calibri',
            'body_font': 'Calibri',
            'fonts_used': set()
        }

        try:
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame'):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.font.name:
                                    fonts['fonts_used'].add(run.font.name)

            fonts['fonts_used'] = list(fonts['fonts_used'])

            # Set primary fonts
            if fonts['fonts_used']:
                fonts['body_font'] = fonts['fonts_used'][0]
                if len(fonts['fonts_used']) > 1:
                    fonts['title_font'] = fonts['fonts_used'][1]
                else:
                    fonts['title_font'] = fonts['fonts_used'][0]

        except Exception as e:
            logger.warning(f"Error extracting fonts: {e}")

        return fonts

    def _extract_images(self, prs: Presentation) -> List[Dict]:
        """Extract and encode images from the presentation"""
        images = []

        try:
            for slide_idx, slide in enumerate(prs.slides):
                for shape_idx, shape in enumerate(slide.shapes):
                    if hasattr(shape, 'image'):
                        try:
                            # Get image data
                            image_data = shape.image.blob

                            # Convert to base64 for storage
                            image_b64 = base64.b64encode(
                                image_data).decode('utf-8')

                            # Get image dimensions and position
                            image_info = {
                                'slide_index': slide_idx,
                                'shape_index': shape_idx,
                                'left': shape.left,
                                'top': shape.top,
                                'width': shape.width,
                                'height': shape.height,
                                'data': image_b64,
                                'content_type': self._get_image_content_type(image_data),
                                'size': len(image_data)
                            }

                            images.append(image_info)

                        except Exception as e:
                            logger.warning(
                                f"Error extracting image from slide {slide_idx}: {e}")

        except Exception as e:
            logger.warning(f"Error extracting images: {e}")

        return images

    def _get_image_content_type(self, image_data: bytes) -> str:
        """Determine image content type from data"""
        if image_data.startswith(b'\x89PNG'):
            return 'image/png'
        elif image_data.startswith(b'\xFF\xD8\xFF'):
            return 'image/jpeg'
        elif image_data.startswith(b'GIF'):
            return 'image/gif'
        elif image_data.startswith(b'BM'):
            return 'image/bmp'
        else:
            return 'image/unknown'

    def _get_slide_size(self, prs: Presentation) -> Dict:
        """Get slide dimensions"""
        return {
            'width': prs.slide_width,
            'height': prs.slide_height,
            'width_inches': prs.slide_width / Inches(1),
            'height_inches': prs.slide_height / Inches(1)
        }

    def _analyze_master_slides(self, prs: Presentation) -> Dict:
        """Analyze slide master information"""
        master_info = {
            'layouts_count': len(prs.slide_layouts),
            'has_master': prs.slide_master is not None,
            'master_shapes': []
        }

        try:
            if prs.slide_master:
                for shape in prs.slide_master.shapes:
                    shape_info = {
                        'type': shape.shape_type,
                        'left': shape.left,
                        'top': shape.top,
                        'width': shape.width,
                        'height': shape.height
                    }

                    if hasattr(shape, 'text_frame') and shape.text_frame.text:
                        shape_info['has_text'] = True
                        shape_info['text_sample'] = shape.text_frame.text[:50]

                    master_info['master_shapes'].append(shape_info)

        except Exception as e:
            logger.warning(f"Error analyzing master slides: {e}")

        return master_info

    def get_best_layout_for_slide_type(self, template_data: Dict, slide_type: str) -> Optional[Dict]:
        """Get the best layout for a given slide type"""
        layouts = template_data.get('layouts', [])

        # Mapping of slide types to layout preferences
        type_mapping = {
            'title': ['title', 'basic'],
            'content': ['content', 'basic'],
            'section': ['title', 'content'],
            'conclusion': ['content', 'basic'],
            'image': ['image', 'content']
        }

        preferred_types = type_mapping.get(slide_type, ['basic', 'content'])

        # Find best matching layout
        for pref_type in preferred_types:
            for layout in layouts:
                if layout.get('used_for') == pref_type:
                    return layout

        # Return first layout if no match found
        return layouts[0] if layouts else None

    def extract_slide_structure(self, template_path: str) -> Dict:
        """Extract structural information for slide generation"""
        try:
            prs = Presentation(template_path)

            structure = {
                'slide_count': len(prs.slides),
                'slide_examples': []
            }

            # Analyze first few slides as examples
            for i, slide in enumerate(prs.slides[:3]):
                slide_info = {
                    'slide_number': i + 1,
                    'layout_name': slide.slide_layout.name,
                    'shapes': [],
                    'text_content': []
                }

                for shape in slide.shapes:
                    shape_data = {
                        'type': shape.shape_type,
                        'has_text': hasattr(shape, 'text_frame') and shape.text_frame.text.strip()
                    }

                    if shape_data['has_text']:
                        slide_info['text_content'].append(
                            shape.text_frame.text[:100])

                    slide_info['shapes'].append(shape_data)

                structure['slide_examples'].append(slide_info)

            return structure

        except Exception as e:
            logger.error(f"Error extracting slide structure: {e}")
            return {'slide_count': 0, 'slide_examples': []}

    def validate_template(self, template_path: str) -> Dict[str, Any]:
        """Validate template file and return compatibility info"""
        validation_result = {
            'valid': False,
            'errors': [],
            'warnings': [],
            'compatibility_score': 0,
            'recommendations': []
        }

        try:
            # Check file exists and is readable
            if not os.path.exists(template_path):
                validation_result['errors'].append("Template file not found")
                return validation_result

            # Try to open presentation
            prs = Presentation(template_path)
            validation_result['valid'] = True

            # Check layouts
            layout_count = len(prs.slide_layouts)
            if layout_count < 2:
                validation_result['warnings'].append(
                    "Template has very few layouts")
            else:
                validation_result['compatibility_score'] += 30

            # Check for master slide
            if prs.slide_master:
                validation_result['compatibility_score'] += 20

            # Check for existing slides (examples)
            slide_count = len(prs.slides)
            if slide_count > 0:
                validation_result['compatibility_score'] += 25
                validation_result['recommendations'].append(
                    "Template contains example slides")

            # Check for images
            image_count = 0
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'image'):
                        image_count += 1

            if image_count > 0:
                validation_result['compatibility_score'] += 25
                validation_result['recommendations'].append(
                    f"Template contains {image_count} reusable images")

            validation_result['summary'] = {
                'layouts': layout_count,
                'slides': slide_count,
                'images': image_count,
                'has_master': prs.slide_master is not None
            }

        except Exception as e:
            validation_result['errors'].append(
                f"Failed to analyze template: {str(e)}")

        return validation_result
